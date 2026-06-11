import os
import re
import json
import shutil
import zipfile
import platform
from datetime import datetime
from typing import List, Dict, Tuple

import pandas as pd
import streamlit as st
from dotenv import load_dotenv
from PIL import Image
import pytesseract
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from pdf2image import convert_from_bytes
from fpdf import FPDF
from rank_bm25 import BM25Okapi

import google.generativeai as genai

from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_core.documents import Document as LCDocument
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate


load_dotenv()
api_key = os.getenv("GOOGLE_API_KEY")

if not api_key:
    st.error("GOOGLE_API_KEY is missing. Add it inside your .env file.")
    st.stop()

genai.configure(api_key=api_key)

if platform.system() == "Windows":
    tesseract_path = r"C:\\Program Files\\Tesseract-OCR\\tesseract.exe"
    if os.path.exists(tesseract_path):
        pytesseract.pytesseract.tesseract_cmd = tesseract_path


COLLECTIONS_DIR = "collections"
BACKUPS_DIR = "backups"
DEFAULT_COLLECTION = "default"
DOCUMENTS_JSON = "documents.json"

SEARCH_CURRENT_COLLECTION = "Current Collection"
SEARCH_ALL_COLLECTIONS = "All Collections"


def apply_custom_css():
    st.markdown(
        """
        <style>
        .stApp {
            background:
                radial-gradient(circle at top left, rgba(147, 197, 253, 0.20), transparent 34%),
                radial-gradient(circle at top right, rgba(196, 181, 253, 0.18), transparent 32%),
                linear-gradient(135deg, #eaf1f8 0%, #e7eef7 45%, #edf0f6 100%);
            color: #172033;
        }

        header[data-testid="stHeader"] { background: transparent; }

        .block-container {
            padding-top: 1.25rem;
            padding-bottom: 2rem;
            max-width: 1180px;
        }

        section[data-testid="stSidebar"] {
            background: linear-gradient(180deg, #f4f8fb 0%, #e7eef6 100%);
            border-right: 2px solid #93c5fd;
            box-shadow: 8px 0 28px rgba(37, 99, 235, 0.10);
        }

        section[data-testid="stSidebar"] * { color: #172033 !important; }

        .main-title {
            font-size: 1.78rem;
            font-weight: 850;
            color: #163b88;
            margin-bottom: 0.28rem;
            letter-spacing: -0.02em;
        }

        .sub-title {
            font-size: 0.96rem;
            color: #52637a;
            line-height: 1.55;
        }

        .hero-card {
            padding: 1.25rem 1.7rem;
            border-radius: 22px;
            background: linear-gradient(135deg, #f9fbfd 0%, #f1f6fb 100%);
            border: 1px solid #cbdcf0;
            box-shadow: 0 10px 28px rgba(30, 64, 175, 0.07);
            margin-bottom: 1.6rem;
        }

        .collection-card {
            padding: 0.95rem;
            border-radius: 16px;
            background: linear-gradient(135deg, #eef6ff 0%, #edf4fb 100%);
            border: 1px solid #bfd4ef;
            margin-bottom: 1rem;
            font-weight: 750;
            color: #0f3f5f !important;
        }

        .metric-card {
            padding: 0.75rem;
            border-radius: 15px;
            background: linear-gradient(135deg, #f8fbfd 0%, #eef5fb 100%);
            border: 1px solid #cbdcf0;
            text-align: center;
            box-shadow: 0 6px 16px rgba(15, 23, 42, 0.035);
            margin-bottom: 0.55rem;
        }

        .metric-number {
            font-size: 1.35rem;
            font-weight: 850;
            color: #2563eb;
        }

        .metric-label {
            color: #52637a;
            font-size: 0.8rem;
            font-weight: 650;
        }

        .status-box {
            padding: 0.95rem;
            border-radius: 16px;
            background: linear-gradient(135deg, #e4f4ee 0%, #e7f0fa 100%);
            border: 1px solid #b8d4e8;
            color: #0f3f5f !important;
            font-weight: 750;
            margin-top: 1rem;
            line-height: 1.55;
            box-shadow: 0 8px 20px rgba(15, 23, 42, 0.035);
        }

        .file-card {
            padding: 0.75rem 0.9rem;
            border-radius: 13px;
            background: linear-gradient(135deg, #f8fbfd 0%, #eef5fb 100%);
            border: 1px solid #cbdcf0;
            margin-bottom: 0.5rem;
            box-shadow: 0 5px 14px rgba(15, 23, 42, 0.03);
        }

        .source-card {
            padding: 0.75rem 0.9rem;
            border-radius: 13px;
            background: #eef6ff;
            border: 1px solid #bfd4ef;
            margin-bottom: 0.55rem;
            font-size: 0.92rem;
        }

        .score-pill {
            display: inline-block;
            padding: 0.18rem 0.55rem;
            border-radius: 999px;
            background: #dbeafe;
            color: #1e3a8a;
            font-size: 0.78rem;
            font-weight: 800;
            margin-left: 0.4rem;
        }

        .collection-pill {
            display: inline-block;
            padding: 0.18rem 0.55rem;
            border-radius: 999px;
            background: #dcfce7;
            color: #166534;
            font-size: 0.78rem;
            font-weight: 800;
            margin-left: 0.4rem;
        }

        .chat-user-card {
            padding: 0.9rem 1rem;
            border-radius: 16px;
            background: linear-gradient(135deg, #dbeafe 0%, #e0f2fe 100%);
            border: 1px solid #bfdbfe;
            margin-bottom: 0.7rem;
        }

        .chat-ai-card {
            padding: 0.95rem 1rem;
            border-radius: 16px;
            background: linear-gradient(135deg, #f8fafc 0%, #edf5fb 100%);
            border: 1px solid #cbdcf0;
            margin-bottom: 1rem;
            box-shadow: 0 8px 18px rgba(15, 23, 42, 0.04);
        }

        .chat-label {
            font-size: 0.82rem;
            font-weight: 800;
            color: #1e3a8a;
            margin-bottom: 0.35rem;
        }

        .small-note {
            color: #64748b !important;
            font-size: 0.86rem;
            line-height: 1.45;
        }

        div[data-testid="stExpander"] {
            background: linear-gradient(180deg, #f9fbfd 0%, #f2f6fa 100%);
            border: 1px solid #cbdcf0;
            border-radius: 18px;
            box-shadow: 0 12px 28px rgba(30, 64, 175, 0.07);
            overflow: hidden;
        }

        div[data-testid="stExpander"] summary {
            background: linear-gradient(90deg, #eaf2fb 0%, #eef1fa 100%);
            color: #163b88 !important;
            font-weight: 800;
            border-bottom: 1px solid #d2dfef;
            padding: 0.78rem 1rem !important;
        }

        div[data-testid="stExpander"] summary p {
            color: #163b88 !important;
            font-weight: 800 !important;
        }

        label { color: #334155 !important; font-weight: 650 !important; }

        .stTextInput input,
        .stTextArea textarea {
            background: #f8fafc !important;
            color: #172033 !important;
            border: 1px solid #aebfd4 !important;
            border-radius: 14px !important;
            box-shadow: inset 0 1px 2px rgba(15, 23, 42, 0.03) !important;
        }

        .stTextInput input:focus,
        .stTextArea textarea:focus {
            border: 1px solid #5b9beb !important;
            box-shadow: 0 0 0 3px rgba(91, 155, 235, 0.16) !important;
        }

        .stTextInput input::placeholder,
        .stTextArea textarea::placeholder { color: #7b8aa0 !important; }

        div[data-baseweb="select"] > div {
            background: #f8fafc !important;
            color: #172033 !important;
            border: 1px solid #aebfd4 !important;
            border-radius: 14px !important;
        }

        div[data-baseweb="select"] span { color: #172033 !important; }

        .stButton button {
            background: linear-gradient(135deg, #2563eb 0%, #0ea5e9 100%) !important;
            color: white !important;
            border: none !important;
            border-radius: 14px !important;
            padding: 0.65rem 1rem !important;
            font-weight: 750 !important;
            box-shadow: 0 8px 18px rgba(14, 165, 233, 0.18);
            transition: all 0.18s ease-in-out;
        }

        .stButton button:hover {
            background: linear-gradient(135deg, #1d4ed8 0%, #0284c7 100%) !important;
            color: white !important;
            transform: translateY(-1px);
            box-shadow: 0 10px 22px rgba(14, 165, 233, 0.24);
        }

        .stDownloadButton button {
            background: linear-gradient(135deg, #059669 0%, #34d399 100%) !important;
            color: white !important;
            border-radius: 14px !important;
            font-weight: 750 !important;
            border: none !important;
        }

        [data-testid="stFileUploader"] {
            background: linear-gradient(135deg, #f6f9fc 0%, #edf4fb 100%);
            border: 1px dashed #8cb9e8;
            border-radius: 18px;
            padding: 0.8rem;
        }

        [data-testid="stFileUploader"] section {
            background: #f8fafc !important;
            border: 1px solid #cbdcf0 !important;
            border-radius: 15px !important;
        }

        [data-testid="stFileUploader"] button {
            background: linear-gradient(135deg, #2563eb 0%, #0ea5e9 100%) !important;
            color: #ffffff !important;
            border-radius: 12px !important;
            border: none !important;
            font-weight: 750 !important;
        }

        .stProgress > div > div > div > div {
            background: linear-gradient(90deg, #2563eb, #0ea5e9) !important;
        }

        .stAlert { border-radius: 14px !important; }

        hr {
            border-color: #cfdced;
            margin-top: 1rem;
            margin-bottom: 1.4rem;
        }

        button[kind="header"],
        [data-testid="collapsedControl"],
        [data-testid="stSidebarCollapsedControl"],
        button[data-testid="collapsedControl"],
        button[data-testid="stSidebarCollapsedControl"] {
            background: linear-gradient(135deg, #1d4ed8, #0284c7) !important;
            color: #ffffff !important;
            border-radius: 14px !important;
            box-shadow: 0 10px 26px rgba(37, 99, 235, 0.45) !important;
            border: 2px solid #ffffff !important;
            opacity: 1 !important;
            z-index: 999999 !important;
        }

        button[kind="header"] svg,
        [data-testid="collapsedControl"] svg,
        [data-testid="stSidebarCollapsedControl"] svg,
        button[data-testid="collapsedControl"] svg,
        button[data-testid="stSidebarCollapsedControl"] svg,
        [data-testid="collapsedControl"] svg path,
        [data-testid="stSidebarCollapsedControl"] svg path {
            color: #ffffff !important;
            stroke: #ffffff !important;
            fill: none !important;
            opacity: 1 !important;
        }

        button[kind="header"]:hover,
        [data-testid="collapsedControl"]:hover,
        [data-testid="stSidebarCollapsedControl"]:hover,
        button[data-testid="collapsedControl"]:hover,
        button[data-testid="stSidebarCollapsedControl"]:hover {
            background: linear-gradient(135deg, #1e40af, #0369a1) !important;
            transform: scale(1.06);
        }

        @media (max-width: 768px) {
            .main-title { font-size: 1.35rem; }
            .hero-card { padding: 1rem; }
            .block-container { padding-left: 1rem; padding-right: 1rem; }
        }
        </style>
        """,
        unsafe_allow_html=True
    )


def initialize_session_state():
    if "chat_histories" not in st.session_state:
        st.session_state.chat_histories = {}
    if "active_collection" not in st.session_state:
        st.session_state.active_collection = DEFAULT_COLLECTION


def sanitize_collection_name(name: str) -> str:
    name = name.strip().lower()
    return re.sub(r"[^a-zA-Z0-9_-]", "_", name)


def ensure_collections_dir():
    os.makedirs(COLLECTIONS_DIR, exist_ok=True)


def ensure_backups_dir():
    os.makedirs(BACKUPS_DIR, exist_ok=True)


def get_collection_path(collection_name: str) -> str:
    ensure_collections_dir()
    return os.path.join(COLLECTIONS_DIR, collection_name)


def get_documents_json_path(collection_name: str) -> str:
    return os.path.join(get_collection_path(collection_name), DOCUMENTS_JSON)


def collection_exists(collection_name: str) -> bool:
    return os.path.exists(get_collection_path(collection_name))


def list_collections():
    ensure_collections_dir()
    collections = []
    for item in os.listdir(COLLECTIONS_DIR):
        full_path = os.path.join(COLLECTIONS_DIR, item)
        if os.path.isdir(full_path):
            collections.append(item)
    collections.sort()
    if DEFAULT_COLLECTION not in collections:
        collections.insert(0, DEFAULT_COLLECTION)
    return collections


def create_collection(collection_name: str):
    collection_name = sanitize_collection_name(collection_name)
    if not collection_name:
        return None
    os.makedirs(get_collection_path(collection_name), exist_ok=True)
    return collection_name


def delete_collection(collection_name: str):
    if collection_name == DEFAULT_COLLECTION:
        return False
    path = get_collection_path(collection_name)
    if os.path.exists(path):
        shutil.rmtree(path)
    if collection_name in st.session_state.chat_histories:
        del st.session_state.chat_histories[collection_name]
    return True


def get_collection_index_path(collection_name: str):
    return get_collection_path(collection_name)


def collection_has_index(collection_name: str):
    path = get_collection_index_path(collection_name)
    return os.path.exists(os.path.join(path, "index.faiss")) and os.path.exists(os.path.join(path, "index.pkl"))


def collection_has_documents_json(collection_name: str):
    return os.path.exists(get_documents_json_path(collection_name))


def reset_collection_index(collection_name: str):
    path = get_collection_index_path(collection_name)
    for file_name in ["index.faiss", "index.pkl", DOCUMENTS_JSON]:
        file_path = os.path.join(path, file_name)
        if os.path.exists(file_path):
            os.remove(file_path)


def get_timestamp_label():
    return datetime.now().strftime("%Y%m%d_%H%M%S")


def get_backup_file_name(prefix: str):
    safe_prefix = sanitize_collection_name(prefix)
    return f"{safe_prefix}_{get_timestamp_label()}.zip"


def zip_folder(source_folder: str, output_zip_path: str):
    with zipfile.ZipFile(output_zip_path, "w", zipfile.ZIP_DEFLATED) as zip_file:
        for root, _, files in os.walk(source_folder):
            for file_name in files:
                file_path = os.path.join(root, file_name)
                arc_name = os.path.relpath(file_path, source_folder)
                zip_file.write(file_path, arc_name)
    return output_zip_path


def export_selected_collection(collection_name: str):
    ensure_backups_dir()
    collection_path = get_collection_path(collection_name)
    if not os.path.exists(collection_path):
        return None
    output_name = get_backup_file_name(collection_name)
    output_path = os.path.join(BACKUPS_DIR, output_name)
    zip_folder(collection_path, output_path)
    return output_path


def export_all_collections():
    ensure_backups_dir()
    ensure_collections_dir()
    if not os.path.exists(COLLECTIONS_DIR):
        return None
    output_name = get_backup_file_name("all_collections_backup")
    output_path = os.path.join(BACKUPS_DIR, output_name)
    zip_folder(COLLECTIONS_DIR, output_path)
    return output_path


def import_collection_zip(uploaded_zip_file, collection_name: str):
    ensure_collections_dir()
    collection_name = sanitize_collection_name(collection_name)
    if not collection_name:
        return False, "Please enter a valid collection name."
    if not uploaded_zip_file.name.lower().endswith(".zip"):
        return False, "Please upload a ZIP file."
    target_path = get_collection_path(collection_name)
    if os.path.exists(target_path):
        shutil.rmtree(target_path)
    os.makedirs(target_path, exist_ok=True)
    try:
        with zipfile.ZipFile(uploaded_zip_file, "r") as zip_ref:
            zip_ref.extractall(target_path)
        return True, f"Collection imported as: {collection_name}"
    except Exception as e:
        if os.path.exists(target_path):
            shutil.rmtree(target_path)
        return False, f"Import failed: {e}"


def clear_all_collections():
    if os.path.exists(COLLECTIONS_DIR):
        shutil.rmtree(COLLECTIONS_DIR)
    ensure_collections_dir()
    st.session_state.chat_histories = {}
    st.session_state.active_collection = DEFAULT_COLLECTION


def read_file_as_bytes(file_path: str):
    with open(file_path, "rb") as file:
        return file.read()


def get_chat_history(collection_name: str):
    if collection_name not in st.session_state.chat_histories:
        st.session_state.chat_histories[collection_name] = []
    return st.session_state.chat_histories[collection_name]


def save_chat_message(collection_name: str, question: str, answer: str, source_docs: List[LCDocument], search_mode: str):
    history = get_chat_history(collection_name)
    history.append({"question": question, "answer": answer, "sources": source_docs, "search_mode": search_mode})


def clear_chat_history(collection_name: str):
    st.session_state.chat_histories[collection_name] = []


def build_source_label(metadata: Dict) -> str:
    file_name = metadata.get("file_name", "Unknown file")
    file_type = metadata.get("file_type", "Unknown")
    page = metadata.get("page")
    chunk = metadata.get("chunk")
    if page:
        return f"{file_name} | {file_type} | Page {page}"
    return f"{file_name} | {file_type} | Chunk {chunk}"


def build_chat_history_text(collection_name: str):
    history = get_chat_history(collection_name)
    lines = []
    for index, item in enumerate(history, start=1):
        lines.append(f"Q{index}: {item['question']}")
        lines.append(f"Search Mode: {item.get('search_mode', SEARCH_CURRENT_COLLECTION)}")
        lines.append(f"A{index}: {item['answer']}")
        lines.append("Sources:")
        seen = set()
        for doc in item["sources"]:
            source_collection = doc.metadata.get("collection", collection_name)
            label = build_source_label(doc.metadata)
            full_label = f"{source_collection} | {label}"
            if full_label not in seen:
                seen.add(full_label)
                lines.append(f"- {full_label}")
        lines.append("")
    return "\n".join(lines)


def get_embeddings():
    return GoogleGenerativeAIEmbeddings(model="models/gemini-embedding-001", google_api_key=api_key)


def tokenize_text(text: str):
    return re.findall(r"\b\w+\b", text.lower())


def serialize_documents(documents: List[LCDocument]):
    return [{"page_content": doc.page_content, "metadata": doc.metadata} for doc in documents]


def deserialize_documents(raw_documents) -> List[LCDocument]:
    return [LCDocument(page_content=item.get("page_content", ""), metadata=item.get("metadata", {})) for item in raw_documents]


def save_documents_json(collection_name: str, documents: List[LCDocument]):
    path = get_collection_path(collection_name)
    os.makedirs(path, exist_ok=True)
    with open(get_documents_json_path(collection_name), "w", encoding="utf-8") as file:
        json.dump(serialize_documents(documents), file, ensure_ascii=False, indent=2)


def load_documents_json(collection_name: str) -> List[LCDocument]:
    json_path = get_documents_json_path(collection_name)
    if not os.path.exists(json_path):
        return []
    with open(json_path, "r", encoding="utf-8") as file:
        raw_documents = json.load(file)
    return deserialize_documents(raw_documents)


def split_text_with_metadata(text: str, base_metadata: Dict) -> List[LCDocument]:
    splitter = RecursiveCharacterTextSplitter(chunk_size=5000, chunk_overlap=500)
    chunks = splitter.split_text(text)
    documents = []
    for index, chunk in enumerate(chunks, start=1):
        metadata = base_metadata.copy()
        metadata["chunk"] = index
        documents.append(LCDocument(page_content=chunk, metadata=metadata))
    return documents


def extract_documents_from_pdfs(pdf_docs) -> List[LCDocument]:
    documents = []
    for pdf in pdf_docs:
        try:
            pdf.seek(0)
            reader = PdfReader(pdf)
            has_normal_text = False
            for page_index, page in enumerate(reader.pages, start=1):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    has_normal_text = True
                    documents.extend(split_text_with_metadata(page_text, {"file_name": pdf.name, "file_type": "PDF", "page": page_index}))
            if not has_normal_text:
                pdf.seek(0)
                images = convert_from_bytes(pdf.read())
                for page_index, image in enumerate(images, start=1):
                    scanned_text = pytesseract.image_to_string(image)
                    if scanned_text.strip():
                        documents.extend(split_text_with_metadata(scanned_text, {"file_name": pdf.name, "file_type": "Scanned PDF OCR", "page": page_index}))
        except Exception as e:
            st.error(f"Error extracting PDF text from {pdf.name}: {e}")
    return documents


def extract_documents_from_images(image_docs) -> List[LCDocument]:
    documents = []
    for image_file in image_docs:
        try:
            image = Image.open(image_file)
            image_text = pytesseract.image_to_string(image)
            if image_text.strip():
                documents.extend(split_text_with_metadata(image_text, {"file_name": image_file.name, "file_type": "Image OCR"}))
        except Exception as e:
            st.error(f"Error extracting image text from {image_file.name}: {e}")
    return documents


def extract_documents_from_docx(docx_docs) -> List[LCDocument]:
    documents = []
    for docx_file in docx_docs:
        try:
            document = Document(docx_file)
            docx_text = ""
            for para in document.paragraphs:
                if para.text.strip():
                    docx_text += para.text + "\n"
            for table in document.tables:
                for row in table.rows:
                    row_values = [cell.text.strip() for cell in row.cells]
                    docx_text += " | ".join(row_values) + "\n"
            if docx_text.strip():
                documents.extend(split_text_with_metadata(docx_text, {"file_name": docx_file.name, "file_type": "DOCX"}))
        except Exception as e:
            st.error(f"Error extracting DOCX text from {docx_file.name}: {e}")
    return documents


def extract_documents_from_txt(txt_docs) -> List[LCDocument]:
    documents = []
    for txt_file in txt_docs:
        try:
            txt_file.seek(0)
            raw_data = txt_file.read()
            try:
                text = raw_data.decode("utf-8")
            except UnicodeDecodeError:
                text = raw_data.decode("latin-1", errors="ignore")
            if text.strip():
                documents.extend(split_text_with_metadata(text, {"file_name": txt_file.name, "file_type": "TXT"}))
        except Exception as e:
            st.error(f"Error extracting TXT text from {txt_file.name}: {e}")
    return documents


def extract_documents_from_csv(csv_docs) -> List[LCDocument]:
    documents = []
    for csv_file in csv_docs:
        try:
            csv_file.seek(0)
            dataframe = pd.read_csv(csv_file)
            csv_text = dataframe.to_csv(index=False)
            if csv_text.strip():
                documents.extend(split_text_with_metadata(csv_text, {"file_name": csv_file.name, "file_type": "CSV"}))
        except Exception as e:
            st.error(f"Error extracting CSV text from {csv_file.name}: {e}")
    return documents


def extract_documents_from_xlsx(xlsx_docs) -> List[LCDocument]:
    documents = []
    for xlsx_file in xlsx_docs:
        try:
            xlsx_file.seek(0)
            sheets = pd.read_excel(xlsx_file, sheet_name=None)
            workbook_text = ""
            for sheet_name, dataframe in sheets.items():
                workbook_text += f"\n\n--- Sheet: {sheet_name} ---\n"
                workbook_text += dataframe.to_csv(index=False)
            if workbook_text.strip():
                documents.extend(split_text_with_metadata(workbook_text, {"file_name": xlsx_file.name, "file_type": "XLSX"}))
        except Exception as e:
            st.error(f"Error extracting XLSX text from {xlsx_file.name}: {e}")
    return documents


def extract_documents_from_pptx(pptx_docs) -> List[LCDocument]:
    documents = []
    for pptx_file in pptx_docs:
        try:
            pptx_file.seek(0)
            presentation = Presentation(pptx_file)
            pptx_text = ""
            for slide_index, slide in enumerate(presentation.slides, start=1):
                slide_text = f"\n\n--- Slide {slide_index} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text.strip() + "\n"
                    if getattr(shape, "has_table", False):
                        for row in shape.table.rows:
                            row_values = [cell.text.strip() for cell in row.cells]
                            slide_text += " | ".join(row_values) + "\n"
                if slide_text.strip():
                    pptx_text += slide_text
            if pptx_text.strip():
                documents.extend(split_text_with_metadata(pptx_text, {"file_name": pptx_file.name, "file_type": "PPTX"}))
        except Exception as e:
            st.error(f"Error extracting PPTX text from {pptx_file.name}: {e}")
    return documents


def create_faiss_vector_store(collection_name: str, documents: List[LCDocument]):
    try:
        embeddings = get_embeddings()
        vector_store = FAISS.from_documents(documents, embedding=embeddings)
        collection_path = get_collection_index_path(collection_name)
        os.makedirs(collection_path, exist_ok=True)
        vector_store.save_local(collection_path)
        save_documents_json(collection_name, documents)
        return vector_store
    except Exception as e:
        st.error(f"Error creating FAISS vector store: {e}")
        return None


def load_faiss_vector_store(collection_name: str):
    try:
        embeddings = get_embeddings()
        collection_path = get_collection_index_path(collection_name)
        return FAISS.load_local(collection_path, embeddings, allow_dangerous_deserialization=True)
    except Exception as e:
        st.error(f"Error loading FAISS vector store: {e}")
        return None


def vector_search_documents(collection_name: str, query: str, k: int = 8) -> List[Tuple[LCDocument, float]]:
    vector_store = load_faiss_vector_store(collection_name)
    if not vector_store:
        return []
    try:
        results = vector_store.similarity_search_with_score(query, k=k)
        normalized_results = []
        for doc, distance in results:
            score = 1 / (1 + float(distance))
            doc.metadata["vector_score"] = round(score, 4)
            normalized_results.append((doc, score))
        return normalized_results
    except Exception as e:
        st.error(f"Vector search failed: {e}")
        return []


def bm25_keyword_search_documents(collection_name: str, query: str, k: int = 8) -> List[Tuple[LCDocument, float]]:
    documents = load_documents_json(collection_name)
    if not documents:
        return []
    tokenized_docs = [tokenize_text(doc.page_content) for doc in documents]
    query_tokens = tokenize_text(query)
    if not query_tokens:
        return []
    bm25 = BM25Okapi(tokenized_docs)
    scores = bm25.get_scores(query_tokens)
    scored_docs = []
    max_score = max(scores) if len(scores) else 0
    for doc, score in zip(documents, scores):
        if score <= 0:
            continue
        normalized_score = float(score) / float(max_score) if max_score > 0 else 0
        doc.metadata["keyword_score"] = round(normalized_score, 4)
        scored_docs.append((doc, normalized_score))
    scored_docs.sort(key=lambda item: item[1], reverse=True)
    return scored_docs[:k]


def make_doc_key(doc: LCDocument) -> str:
    metadata = doc.metadata
    return "|".join([metadata.get("collection", ""), metadata.get("file_name", ""), metadata.get("file_type", ""), str(metadata.get("page", "")), str(metadata.get("chunk", ""))])


def hybrid_search_documents(collection_name: str, query: str, k: int = 8) -> List[LCDocument]:
    vector_results = vector_search_documents(collection_name, query, k=k)
    keyword_results = bm25_keyword_search_documents(collection_name, query, k=k)
    merged = {}
    for doc, score in vector_results:
        doc.metadata["collection"] = collection_name
        key = make_doc_key(doc)
        if key not in merged:
            merged[key] = {"doc": doc, "vector_score": 0.0, "keyword_score": 0.0}
        merged[key]["vector_score"] = max(merged[key]["vector_score"], score)
    for doc, score in keyword_results:
        doc.metadata["collection"] = collection_name
        key = make_doc_key(doc)
        if key not in merged:
            merged[key] = {"doc": doc, "vector_score": 0.0, "keyword_score": 0.0}
        merged[key]["keyword_score"] = max(merged[key]["keyword_score"], score)
    ranked = []
    for item in merged.values():
        vector_score = item["vector_score"]
        keyword_score = item["keyword_score"]
        final_score = (vector_score * 0.65) + (keyword_score * 0.35)
        doc = item["doc"]
        doc.metadata["vector_score"] = round(vector_score, 4)
        doc.metadata["keyword_score"] = round(keyword_score, 4)
        doc.metadata["final_score"] = round(final_score, 4)
        ranked.append(doc)
    ranked.sort(key=lambda doc: doc.metadata.get("final_score", 0), reverse=True)
    return ranked[:k]


def get_searchable_collections():
    return [collection for collection in list_collections() if collection_has_index(collection) and collection_has_documents_json(collection)]


def search_single_collection(collection_name: str, query: str, k: int = 8):
    return hybrid_search_documents(collection_name=collection_name, query=query, k=k)


def search_all_collections(query: str, k: int = 10):
    all_results = []
    for collection_name in get_searchable_collections():
        try:
            all_results.extend(search_single_collection(collection_name=collection_name, query=query, k=k))
        except Exception as e:
            st.warning(f"Search skipped for {collection_name}: {e}")
    all_results.sort(key=lambda doc: doc.metadata.get("final_score", 0), reverse=True)
    return all_results[:k]


def search_documents(collection_name: str, query: str, search_mode: str, k: int = 8):
    if search_mode == SEARCH_ALL_COLLECTIONS:
        return search_all_collections(query=query, k=k)
    return search_single_collection(collection_name=collection_name, query=query, k=k)


def get_total_collection_count():
    return len(get_searchable_collections())


def get_total_index_count():
    total_chunks = 0
    for collection_name in get_searchable_collections():
        total_chunks += len(load_documents_json(collection_name))
    return total_chunks


def build_collection_summary():
    return {"collections": get_total_collection_count(), "chunks": get_total_index_count()}


def render_search_mode_card(search_mode: str):
    if search_mode == SEARCH_ALL_COLLECTIONS:
        stats = build_collection_summary()
        st.markdown(
            f"""
            <div class="collection-card">
                🌎 Search Scope: <b>All Collections</b><br>
                <span class="small-note">
                Collections: {stats['collections']} |
                Chunks: {stats['chunks']}
                </span>
            </div>
            """,
            unsafe_allow_html=True
        )
    else:
        st.markdown(
            """
            <div class="collection-card">
                📁 Search Scope: <b>Current Collection</b>
            </div>
            """,
            unsafe_allow_html=True
        )


def show_sources(docs: List[LCDocument], title: str = "Sources"):
    if not docs:
        return
    st.markdown(f"### 🔎 {title}")
    seen = set()
    unique_docs = []
    for doc in docs:
        label = build_source_label(doc.metadata)
        collection = doc.metadata.get("collection", "unknown")
        full_label = f"{collection} | {label}"
        if full_label not in seen:
            seen.add(full_label)
            unique_docs.append(doc)
    for i, doc in enumerate(unique_docs, start=1):
        label = build_source_label(doc.metadata)
        collection = doc.metadata.get("collection", "unknown")
        final_score = doc.metadata.get("final_score", 0)
        vector_score = doc.metadata.get("vector_score", 0)
        keyword_score = doc.metadata.get("keyword_score", 0)
        preview = doc.page_content[:900].strip()
        st.markdown(
            f"""
            <div class="source-card">
                <b>{i}. {label}</b>
                <span class="collection-pill">{collection}</span>
                <span class="score-pill">Score: {final_score}</span>
                <br>
                <span class="small-note">
                    Vector: {vector_score} · Keyword: {keyword_score}
                </span>
            </div>
            """,
            unsafe_allow_html=True
        )
        with st.expander(f"View source preview {i}", expanded=False):
            st.write(preview)


def initialize_qa_chain():
    prompt_template = """
You are a helpful document assistant.

Answer the user's question using only the provided context.

Rules:
- Be clear and direct.
- Do not invent information.
- If the answer is not available in the context, say:
  "I could not find this information in the uploaded document."
- The app will show source citations separately below your answer.
- The context was selected using hybrid search with vector and keyword ranking.

Context:
{context}

Question:
{question}

Answer:
"""
    model = ChatGoogleGenerativeAI(model="gemini-2.5-flash", temperature=0.2, google_api_key=api_key)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    return load_qa_chain(model, chain_type="stuff", prompt=prompt)


def initialize_summary_chain():
    prompt_template = """
You are a document summarization assistant.

Use only the provided context.

Instruction:
{instruction}

Summary Length:
{summary_length_instruction}

Rules:
- Summarize clearly.
- Do not invent information.
- Keep the result useful and structured.
- The app will show source citations separately below the summary.
- The context was selected using hybrid search with vector and keyword ranking.

Context:
{context}

Write the summary:
"""
    model = ChatGoogleGenerativeAI(model="gemini-2.5-flash", temperature=0.3, google_api_key=api_key)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "instruction", "summary_length_instruction"])
    return load_qa_chain(model, chain_type="stuff", prompt=prompt)


def summarize_documents(collection_name: str, user_instruction, search_mode: str, topic=None, summary_length="short") -> Tuple[str, List[LCDocument]]:
    query = user_instruction
    if topic:
        query += " " + topic
    docs = search_documents(collection_name=collection_name, query=query, search_mode=search_mode, k=8)
    if not docs:
        return "", []
    length_mapping = {
        "short": "Write a short summary with only the most important points.",
        "medium": "Write a medium-length summary with key points and useful details.",
        "long": "Write a detailed summary with important explanations and structure."
    }
    try:
        chain = initialize_summary_chain()
        response = chain(
            {
                "input_documents": docs,
                "instruction": user_instruction,
                "summary_length_instruction": length_mapping.get(summary_length, length_mapping["short"])
            },
            return_only_outputs=True
        )
        return response["output_text"], docs
    except Exception as e:
        st.error(f"Error generating summary: {e}")
        return "", docs


def answer_user_question(collection_name: str, user_question, search_mode: str, topic=None) -> Tuple[str, List[LCDocument]]:
    query = user_question
    if topic:
        query += " " + topic
    docs = search_documents(collection_name=collection_name, query=query, search_mode=search_mode, k=8)
    if not docs:
        return "", []
    try:
        chain = initialize_qa_chain()
        response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)
        return response["output_text"], docs
    except Exception as e:
        st.error(f"Error generating answer: {e}")
        return "", docs


def generate_pdf_summary(summary):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    safe_summary = summary.encode("latin-1", "replace").decode("latin-1")
    pdf.multi_cell(0, 10, txt=safe_summary, align="L")
    return pdf.output(dest="S").encode("latin1")


def download_pdf(summary):
    pdf_data = generate_pdf_summary(summary)
    st.download_button(label="⬇️ Download PDF Summary", data=pdf_data, file_name="summary.pdf", mime="application/pdf")


def download_chat_history(collection_name: str):
    history = get_chat_history(collection_name)
    if not history:
        return
    chat_text = build_chat_history_text(collection_name)
    st.download_button(label="⬇️ Download Chat History", data=chat_text, file_name=f"{collection_name}_chat_history.txt", mime="text/plain")


def render_chat_history(collection_name: str):
    history = get_chat_history(collection_name)
    if not history:
        st.info("No chat yet. Ask a question after processing your selected collection.")
        return
    for index, item in enumerate(history, start=1):
        search_mode = item.get("search_mode", SEARCH_CURRENT_COLLECTION)
        st.markdown(
            f"""
            <div class="chat-user-card">
                <div class="chat-label">
                    You · Question {index} · {search_mode}
                </div>
                {item["question"]}
            </div>
            """,
            unsafe_allow_html=True
        )
        st.markdown(
            f"""
            <div class="chat-ai-card">
                <div class="chat-label">Assistant · Answer {index}</div>
                {item["answer"]}
            </div>
            """,
            unsafe_allow_html=True
        )
        show_sources(item["sources"], title=f"Sources for Answer {index}")


def get_file_type_counts(uploaded_files):
    counts = {"PDF": 0, "Image": 0, "DOCX": 0, "TXT": 0, "CSV": 0, "XLSX": 0, "PPTX": 0}
    for file in uploaded_files:
        name = file.name.lower()
        if name.endswith(".pdf"):
            counts["PDF"] += 1
        elif name.endswith((".jpg", ".jpeg", ".png")):
            counts["Image"] += 1
        elif name.endswith(".docx"):
            counts["DOCX"] += 1
        elif name.endswith(".txt"):
            counts["TXT"] += 1
        elif name.endswith(".csv"):
            counts["CSV"] += 1
        elif name.endswith(".xlsx"):
            counts["XLSX"] += 1
        elif name.endswith(".pptx"):
            counts["PPTX"] += 1
    return counts


def main():
    st.set_page_config(page_title="Document Summary Assistant", page_icon="📄", layout="wide")
    initialize_session_state()
    apply_custom_css()
    ensure_collections_dir()
    ensure_backups_dir()

    st.markdown(
        """
        <div class="hero-card">
            <div class="main-title">📄 Document Summary Assistant</div>
            <div class="sub-title">
                Search one collection or all collections using hybrid vector + keyword retrieval,
                source citations, relevance scores, chat history, local backups, and expanded file support.
            </div>
        </div>
        """,
        unsafe_allow_html=True
    )

    with st.sidebar:
        st.markdown("## 🗂️ Document Library")

        new_collection_name = st.text_input("Create collection", placeholder="Example: project_report")

        if st.button("➕ Create Collection", use_container_width=True):
            created_name = create_collection(new_collection_name)
            if created_name:
                st.session_state.active_collection = created_name
                st.success(f"Collection created: {created_name}")
            else:
                st.warning("Please enter a valid collection name.")

        collections = list_collections()
        selected_collection = st.selectbox(
            "Active collection",
            options=collections,
            index=collections.index(st.session_state.active_collection) if st.session_state.active_collection in collections else 0
        )
        st.session_state.active_collection = selected_collection
        active_collection = st.session_state.active_collection

        search_mode = st.radio("Search mode", options=[SEARCH_CURRENT_COLLECTION, SEARCH_ALL_COLLECTIONS], horizontal=False)

        st.markdown(
            f"""
            <div class="collection-card">
                📁 Active Collection:<br>
                <b>{active_collection}</b><br>
                <span class="small-note">Mode: {search_mode}</span>
            </div>
            """,
            unsafe_allow_html=True
        )

        st.markdown("## 💾 Local Backup")
        st.markdown(
            """
            <div class="status-box">
                🖥️ Mode: Local Private Workspace<br>
                <span class="small-note">
                    Data is stored only on this computer inside collections/.
                </span>
            </div>
            """,
            unsafe_allow_html=True
        )

        selected_export_path = export_selected_collection(active_collection)
        if selected_export_path:
            st.download_button(
                label="⬇️ Export Selected Collection",
                data=read_file_as_bytes(selected_export_path),
                file_name=os.path.basename(selected_export_path),
                mime="application/zip",
                use_container_width=True
            )

        all_export_path = export_all_collections()
        if all_export_path:
            st.download_button(
                label="⬇️ Export All Collections",
                data=read_file_as_bytes(all_export_path),
                file_name=os.path.basename(all_export_path),
                mime="application/zip",
                use_container_width=True
            )

        with st.expander("📥 Import Collection ZIP", expanded=False):
            import_name = st.text_input("Import as collection name", placeholder="Example: imported_project")
            import_zip = st.file_uploader("Upload collection ZIP", type=["zip"], key="collection_zip_import")

            if st.button("📥 Import ZIP", use_container_width=True):
                if import_zip is None:
                    st.warning("Please upload a ZIP file.")
                else:
                    success, message = import_collection_zip(import_zip, import_name)
                    if success:
                        st.session_state.active_collection = sanitize_collection_name(import_name)
                        st.success(message)
                    else:
                        st.error(message)

        with st.expander("⚠️ Danger Zone", expanded=False):
            st.warning("This removes all local collections from this computer.")
            confirm_clear = st.text_input("Type CLEAR to confirm", key="clear_all_confirm")

            if st.button("🧨 Clear All Local Collections", use_container_width=True):
                if confirm_clear == "CLEAR":
                    clear_all_collections()
                    st.success("All local collections cleared.")
                else:
                    st.error("Confirmation text did not match.")

        st.markdown("## 📂 Upload Center")
        st.markdown(
            '<p class="small-note">Supported: PDF, scanned PDF, JPG, PNG, DOCX, TXT, CSV, XLSX, PPTX.</p>',
            unsafe_allow_html=True
        )

        uploaded_files = st.file_uploader(
            "Upload your files",
            accept_multiple_files=True,
            type=["pdf", "jpg", "jpeg", "png", "docx", "txt", "csv", "xlsx", "pptx"]
        )

        if uploaded_files:
            counts = get_file_type_counts(uploaded_files)
            st.markdown("### 📊 Upload Summary")
            metric_cols = st.columns(4)
            file_labels = [
                ("PDF", counts["PDF"]), ("Image", counts["Image"]), ("DOCX", counts["DOCX"]), ("TXT", counts["TXT"]),
                ("CSV", counts["CSV"]), ("XLSX", counts["XLSX"]), ("PPTX", counts["PPTX"])
            ]
            for index, (label, value) in enumerate(file_labels):
                with metric_cols[index % 4]:
                    st.markdown(
                        f"""
                        <div class="metric-card">
                            <div class="metric-number">{value}</div>
                            <div class="metric-label">{label}</div>
                        </div>
                        """,
                        unsafe_allow_html=True
                    )

            st.markdown("### 📁 Selected Files")
            for file in uploaded_files:
                size_kb = round(file.size / 1024, 2)
                st.markdown(
                    f"""
                    <div class="file-card">
                        <b>{file.name}</b><br>
                        <span class="small-note">{size_kb} KB</span>
                    </div>
                    """,
                    unsafe_allow_html=True
                )

        process_button = st.button("🚀 Process Files Into Collection", use_container_width=True)

        if st.button("🧹 Reset Selected Collection Index", use_container_width=True):
            reset_collection_index(active_collection)
            clear_chat_history(active_collection)
            st.success(f"Index and chat reset for: {active_collection}")

        if active_collection != DEFAULT_COLLECTION:
            if st.button("🗑️ Delete Selected Collection", use_container_width=True):
                deleted = delete_collection(active_collection)
                if deleted:
                    st.session_state.active_collection = DEFAULT_COLLECTION
                    st.success("Collection deleted.")
                else:
                    st.warning("Default collection cannot be deleted.")

        if st.button("💬 Clear Collection Chat", use_container_width=True):
            clear_chat_history(active_collection)
            st.success("Chat history cleared for this collection.")

        download_chat_history(active_collection)

        if collection_has_index(active_collection) and collection_has_documents_json(active_collection):
            st.markdown("""<div class="status-box">✅ Active collection is ready.</div>""", unsafe_allow_html=True)
        else:
            st.markdown("""<div class="status-box">ℹ️ No searchable index for active collection. Upload and process files first.</div>""", unsafe_allow_html=True)

        if search_mode == SEARCH_ALL_COLLECTIONS:
            stats = build_collection_summary()
            st.markdown(
                f"""
                <div class="status-box">
                    🌎 All-collection search ready for:<br>
                    Collections: <b>{stats["collections"]}</b><br>
                    Indexed chunks: <b>{stats["chunks"]}</b>
                </div>
                """,
                unsafe_allow_html=True
            )

    if uploaded_files:
        file_names = [file.name for file in uploaded_files]
        st.markdown("### 🧾 File Selection")
        selected_files = st.multiselect("Choose which files to process into active collection", options=file_names, default=file_names)

        if process_button:
            selected_uploaded_files = [file for file in uploaded_files if file.name in selected_files]

            if not selected_uploaded_files:
                st.warning("Please select at least one file.")
                return

            progress_bar = st.progress(0)
            status_area = st.empty()

            with st.spinner(f"Processing files into collection: {active_collection}"):
                status_area.info("Step 1/5: Sorting uploaded files...")
                progress_bar.progress(12)

                pdf_files = [file for file in selected_uploaded_files if file.type == "application/pdf" or file.name.lower().endswith(".pdf")]
                image_files = [file for file in selected_uploaded_files if file.type.startswith("image/")]
                docx_files = [file for file in selected_uploaded_files if file.name.lower().endswith(".docx")]
                txt_files = [file for file in selected_uploaded_files if file.name.lower().endswith(".txt")]
                csv_files = [file for file in selected_uploaded_files if file.name.lower().endswith(".csv")]
                xlsx_files = [file for file in selected_uploaded_files if file.name.lower().endswith(".xlsx")]
                pptx_files = [file for file in selected_uploaded_files if file.name.lower().endswith(".pptx")]

                status_area.info("Step 2/5: Extracting text with metadata...")
                progress_bar.progress(35)

                documents = []
                documents.extend(extract_documents_from_pdfs(pdf_files))
                documents.extend(extract_documents_from_images(image_files))
                documents.extend(extract_documents_from_docx(docx_files))
                documents.extend(extract_documents_from_txt(txt_files))
                documents.extend(extract_documents_from_csv(csv_files))
                documents.extend(extract_documents_from_xlsx(xlsx_files))
                documents.extend(extract_documents_from_pptx(pptx_files))

                if documents:
                    status_area.info("Step 3/5: Saving chunks for keyword search...")
                    progress_bar.progress(58)
                    save_documents_json(active_collection, documents)

                    status_area.info("Step 4/5: Creating vector index...")
                    progress_bar.progress(80)

                    vector_store = create_faiss_vector_store(active_collection, documents)

                    if vector_store:
                        progress_bar.progress(100)
                        status_area.success(f"Collection '{active_collection}' updated with {len(documents)} chunks. Hybrid search is ready.")
                    else:
                        status_area.error("Text was extracted, but vector index creation failed. Check your Google API key.")
                else:
                    status_area.error("No extractable text found.")

    st.divider()
    render_search_mode_card(search_mode)

    st.markdown(
        f"""
        <div class="collection-card">
            Active collection: <b>{active_collection}</b><br>
            <span class="small-note">Search mode: {search_mode}</span>
        </div>
        """,
        unsafe_allow_html=True
    )

    col1, col2 = st.columns(2, gap="large")

    with col1:
        with st.expander("📜 Summarize Documents", expanded=True):
            user_topic = st.text_input("Topic optional", placeholder="Example: security, introduction, conclusion")
            user_instruction = st.text_input("Summary instruction", value="Create a clear summary")
            summary_length = st.selectbox("Summary length", ["short", "medium", "long"])

            if st.button("✍️ Generate Summary", use_container_width=True):
                if search_mode == SEARCH_CURRENT_COLLECTION and not collection_has_index(active_collection):
                    st.warning("Please process files into this collection first.")
                elif search_mode == SEARCH_CURRENT_COLLECTION and not collection_has_documents_json(active_collection):
                    st.warning("Please reprocess this collection to enable hybrid search.")
                elif search_mode == SEARCH_ALL_COLLECTIONS and not get_searchable_collections():
                    st.warning("No searchable collections found. Process files into at least one collection.")
                elif not user_instruction.strip():
                    st.warning("Please enter summary instruction.")
                else:
                    with st.spinner("Generating summary using selected search scope..."):
                        summary, source_docs = summarize_documents(
                            collection_name=active_collection,
                            user_instruction=user_instruction,
                            search_mode=search_mode,
                            topic=user_topic,
                            summary_length=summary_length
                        )
                        if summary:
                            st.success("Summary generated")
                            st.write(summary)
                            download_pdf(summary)
                            show_sources(source_docs, title="Summary Sources")

    with col2:
        with st.expander("💬 Ask Documents", expanded=True):
            user_topic_for_question = st.text_input("Question topic optional", placeholder="Example: eligibility, cost, features")
            user_question = st.text_area("Your question", placeholder="Ask anything from selected search scope...", height=120)

            if st.button("🔍 Ask and Save to Chat", use_container_width=True):
                if search_mode == SEARCH_CURRENT_COLLECTION and not collection_has_index(active_collection):
                    st.warning("Please process files into this collection first.")
                elif search_mode == SEARCH_CURRENT_COLLECTION and not collection_has_documents_json(active_collection):
                    st.warning("Please reprocess this collection to enable hybrid search.")
                elif search_mode == SEARCH_ALL_COLLECTIONS and not get_searchable_collections():
                    st.warning("No searchable collections found. Process files into at least one collection.")
                elif not user_question.strip():
                    st.warning("Please enter a question.")
                else:
                    with st.spinner("Finding answer using selected search scope..."):
                        answer, source_docs = answer_user_question(
                            collection_name=active_collection,
                            user_question=user_question,
                            search_mode=search_mode,
                            topic=user_topic_for_question
                        )
                        if answer:
                            save_chat_message(active_collection, user_question, answer, source_docs, search_mode)
                            st.success("Answer saved to chat history")

    st.divider()
    st.markdown(f"## 💬 Chat History: `{active_collection}`")
    render_chat_history(active_collection)


if __name__ == "__main__":
    main()
